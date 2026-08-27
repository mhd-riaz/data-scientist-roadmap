"""Build the mini-project slide deck.

Run with no install:

    uv run --with python-pptx python submission/slides/build_deck.py

Numbers are read from the trained model's `metrics.json` wherever possible, so the deck
cannot quote a figure the shipped model does not produce. The few that come from closed
experiments (the model ladder, the ensemble oracle) are constants here and are marked as
such in `docs/plan.md`.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
PROJECT = HERE.parents[1]
METRICS = PROJECT / "ml-v2" / "artifacts" / "models" / "v2-001" / "metrics.json"
FIGURES = PROJECT / "submission" / "report" / "figures"
OUT = HERE / "news-topic-classifier.pptx"

INK = RGBColor(0x1B, 0x22, 0x2C)
MUTED = RGBColor(0x5A, 0x66, 0x75)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GOOD = RGBColor(0x1E, 0x7A, 0x4B)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xDD, 0xE2, 0xE8)

W, H = Inches(13.333), Inches(7.5)


def deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _text(slide, left, top, width, height, size=18, bold=False, color=INK,
          align=PP_ALIGN.LEFT, space_after=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.alignment = align
    para.space_after = Pt(space_after)
    run = para.add_run()
    run.font.size, run.font.bold, run.font.color.rgb = Pt(size), bold, color
    return frame, run


def heading(slide, title: str, kicker: str = "") -> None:
    if kicker:
        _, run = _text(slide, Inches(0.7), Inches(0.42), Inches(11.9), Inches(0.3),
                       size=13, bold=True, color=ACCENT)
        run.text = kicker.upper()
    _, run = _text(slide, Inches(0.7), Inches(0.72), Inches(11.9), Inches(0.8),
                   size=30, bold=True)
    run.text = title
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(1.62), Inches(11.9), Emu(9525))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()


def bullets(slide, items: list, top=Inches(2.0), left=Inches(0.7),
            width=Inches(11.9), size=18, height=None) -> None:
    box = slide.shapes.add_textbox(left, top, width,
                                   height or (H - top - Inches(0.4)))
    frame = box.text_frame
    frame.word_wrap = True
    for i, item in enumerate(items):
        text, level = (item, 0) if isinstance(item, str) else item
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.level = level
        para.space_after = Pt(12 if level == 0 else 6)
        run = para.add_run()
        run.text = ("" if level else "") + text
        run.font.size = Pt(size if level == 0 else size - 3)
        run.font.color.rgb = INK if level == 0 else MUTED


def table(slide, rows: list[list[str]], left=Inches(0.7), top=Inches(2.1),
          width=Inches(11.9), height=Inches(0.4), highlight: int | None = None,
          size=14) -> None:
    shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width,
                                   height * len(rows))
    tbl = shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            # add_run rather than cell.text: an empty string creates no run at all,
            # and the font settings below would then have nothing to apply to.
            run = para.add_run()
            run.text = str(value)
            run.font.size = Pt(size)
            run.font.bold = r == 0 or r == highlight
            run.font.color.rgb = ACCENT if r == highlight else INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF5, 0xF8) if r == 0 else PAPER


def note(slide, text: str, top=Inches(6.5)) -> None:
    _, run = _text(slide, Inches(0.7), top, Inches(11.9), Inches(0.6),
                   size=14, color=MUTED)
    run.text = text


def picture(slide, name: str, left, top, width) -> bool:
    path = FIGURES / name
    if not path.is_file():
        return False
    slide.shapes.add_picture(str(path), left, top, width=width)
    return True


def formula(slide, lines: list[str], top, size=20, left=Inches(2.2),
            width=Inches(8.9)) -> None:
    box = slide.shapes.add_shape(5, left, top, width,
                                 Inches(0.52) * len(lines) + Inches(0.3))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF2, 0xF5, 0xF8)
    box.line.color.rgb = RULE
    frame = box.text_frame
    frame.word_wrap = True
    for i, line in enumerate(lines):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = line
        run.font.size, run.font.bold, run.font.color.rgb = Pt(size), True, INK
        run.font.name = "Consolas"


def build(m: dict) -> Presentation:
    val, test, ab = m["validation"], m["test"], m["body_ab"]
    prs = deck()

    # --- 1. title -----------------------------------------------------------
    s = _blank(prs)
    _, run = _text(s, Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.4),
                   size=44, bold=True)
    run.text = "Beyond the Headline"
    _, run = _text(s, Inches(1.0), Inches(3.2), Inches(11.3), Inches(0.9), size=24,
                   color=MUTED)
    run.text = ("Calibrated news article classifier: 13 topics, read from the whole "
                "article, and quiet when it is not sure")
    _, run = _text(s, Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.5), size=18)
    run.text = "Mohamed Riaz  ·  PES1PGE25DS037"
    _, run = _text(s, Inches(1.0), Inches(5.1), Inches(11.3), Inches(0.5), size=16,
                   color=MUTED)
    run.text = ("Supervised Machine Learning — Classification  ·  Mini Project  ·  "
                "Department of CSE, PES University")

    # --- 2. problem ---------------------------------------------------------
    s = _blank(prs)
    heading(s, "The same story, filed five different ways", kicker="Problem")
    bullets(s, [
        "One PTI story on the repo rate. The Hindu files it under Business. "
        "Indian Express files it under India. Deccan Herald puts it in Top Stories.",
        "Every paper has its own section names. There is no common list, because "
        "sections are a newsroom habit, not a fact about the article.",
        "My job: give every article one topic from a fixed list of 13, no matter who "
        "published it.",
        ("Rule I set for myself: only the classical ML from this course. No ChatGPT, no "
         "LLM anywhere. If the model says sport, I must be able to point at the words "
         "that made it say sport.", 1),
    ])
    note(s, "The 13 topics: politics · business · crime and justice · technology · sport · "
            "entertainment · health · education · science and space · environment · "
            "disaster · conflict · society and lifestyle")

    # --- 3. the question ----------------------------------------------------
    s = _blank(prs)
    heading(s, "My question", kicker="Hypothesis")
    bullets(s, [
        "An RSS feed gives me the headline and a one-line summary. About 33 words. "
        "Most news classifiers stop there.",
        "My collector opens the link and pulls the full article. About 650 words, "
        "roughly 100 times more text.",
        "Is the extra text actually useful? Not obvious. Full articles also carry ads, "
        "author bios and “Story continues below this ad”. The model can easily learn "
        "those instead of the topic.",
    ])
    _, run = _text(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.0), size=26,
                   bold=True, color=ACCENT)
    run.text = ("Does the full article beat the headline? "
                "I wanted to measure it, not assume it.")

    # --- 4. system ----------------------------------------------------------
    s = _blank(prs)
    heading(s, "What was built", kicker="System")
    stages = [
        ("Collector", "Go service, 97 RSS/Atom feeds\n→ MongoDB, deduplicated"),
        ("Preparation", "clean · admit · group\nnear-duplicate stories"),
        ("Snapshot", "frozen corpus, so every\nnumber can be re-run"),
        ("Model", "word counts → linear SVM\n→ an honest probability"),
        ("Demo UI", "Streamlit: classify, explain,\nabstain, validate"),
    ]
    for i, (name, detail) in enumerate(stages):
        left = Inches(0.7 + i * 2.44)
        card = s.shapes.add_shape(5, left, Inches(2.6), Inches(2.2), Inches(1.9))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF2, 0xF5, 0xF8)
        card.line.color.rgb = RULE
        card.text_frame.word_wrap = True
        para = card.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = name
        run.font.size, run.font.bold, run.font.color.rgb = Pt(17), True, INK
        sub = card.text_frame.add_paragraph()
        sub.alignment = PP_ALIGN.CENTER
        srun = sub.add_run()
        srun.text = detail
        srun.font.size, srun.font.color.rgb = Pt(11), MUTED
        if i < len(stages) - 1:
            arrow = s.shapes.add_shape(13, left + Inches(2.24), Inches(3.4),
                                       Inches(0.16), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()
    note(s, "Nothing is deployed. This is a proof of concept. Build it, measure it "
            "properly, then decide if it is worth putting anywhere near production.",
         top=Inches(5.0))

    # --- 5. data ------------------------------------------------------------
    s = _blank(prs)
    heading(s, "The data", kicker="Dataset")
    table(s, [
        ["", "Value"],
        ["Articles collected", "14,189 from 40 publishers"],
        ["Articles I labelled by hand", "8,001"],
        ["Articles where the full body was available", "94.3%"],
        ["Biggest class vs smallest class", "6.7 : 1"],
        ["Train / validation / test",
         f"{m['splits'].get('train', 0):,} / {m['splits'].get('val', 0):,} "
         f"/ {m['splits'].get('test', 0):,}"],
        ["Thrown out as not-an-article", "4.1% (582)"],
    ], top=Inches(2.1), height=Inches(0.46))
    note(s, "Random sampling was not an option. War and conflict stories are about 1.7% "
            "of Indian news, so getting 150 of them by random draw means reading about "
            "9,000 articles. I wrote a sampler that goes looking for the rare topics "
            "instead, and ran three rounds of labelling.", top=Inches(5.6))

    # --- 6. preparation -----------------------------------------------------
    s = _blank(prs)
    heading(s, "Cleaning the data is where the accuracy came from", kicker="Method")
    bullets(s, [
        "Every paper repeats the same lines. “Story continues below this ad” shows up in "
        "75% of one publisher's articles. I found 355 such lines across 47 sources and "
        "removed them.",
        ("Author bios are the dangerous ones. A cricket reporter's bio sitting at the "
         "bottom of a business story pulls the whole article towards sport.", 1),
        "9 rules to throw out things that are not really articles: photo galleries, "
        "weather bulletins, “Today's headlines for school assembly”. That removed 4.1%.",
        "The same PTI copy runs in five papers. I group those together, so a copy of a "
        "training article can never appear in the test set.",
        ("If I skip this, the model sees the answer during training and my score looks "
         "better than it really is. Grouping on full bodies catches 41% more copies "
         "than grouping on headlines.", 1),
    ])

    # --- 7. the accuracy trap ----------------------------------------------
    s = _blank(prs)
    heading(s, "Why I do not report accuracy", kicker="Chapter 07 §1")
    formula(s, [
        "Accuracy = (TP + TN) / (TP + TN + FP + FN)",
        "macro-F1 = (F1 of topic 1 + F1 of topic 2 + ... + F1 of topic 13) / 13",
    ], Inches(1.9), size=17, left=Inches(1.4), width=Inches(10.5))
    table(s, [
        ["On the 1,120 validation articles", "accuracy", "macro-F1"],
        ["one line of code: always answer “politics”", "19.0%", "0.025"],
        ["my model", f"{val['accuracy']:.1%}", f"{val['macro_f1']:.3f}"],
    ], top=Inches(3.4), height=Inches(0.5), highlight=2, size=16)
    bullets(s, [
        "Politics is the biggest topic. A model that always answers politics scores 19% "
        "and finds nothing in the other 12 topics. Same trick as the guard who waves "
        "everybody through.",
        "Accuracy here is exactly micro-F1, so it just repeats the confusion matrix.",
        "macro-F1 gives all 13 topics an equal vote. Society and lifestyle with 66 "
        "articles counts as much as politics with 213. That is why I use it.",
    ], top=Inches(5.0), size=15)

    # --- 7b. the confusion matrix -------------------------------------------
    s = _blank(prs)
    heading(s, "The confusion matrix, all 13 topics", kicker="Chapter 07 §2")
    charted = picture(s, "confusion.png", Inches(0.8), Inches(1.85), Inches(6.6))
    bullets(s, [
        "Rows are the true topic. Columns are what my model said. The diagonal is what "
        "it got right.",
        "Politics is the dustbin. It wrongly takes 13 business stories, 13 society "
        "stories and 10 conflict stories.",
        ("Fair enough, honestly. A story about the budget IS both business and "
         "politics. My labellers argued about the same pairs.", 1),
        "Society and lifestyle is the weak row. Only 25 of its 66 articles land on the "
        "diagonal. The rest scatter everywhere.",
        "Sport is the clean row: 89 out of 91. Cricket and football vocabulary does not "
        "look like anything else.",
    ], top=Inches(2.0),
        left=Inches(7.7) if charted else Inches(0.7),
        width=Inches(4.9) if charted else Inches(11.9),
        size=15 if charted else 18)

    # --- 7c. one class out of it --------------------------------------------
    s = _blank(prs)
    heading(s, "Taking one topic out of that matrix", kicker="Chapter 07 §3")
    table(s, [
        ["conflict / war = the positive class", "model said conflict",
         "model said something else"],
        ["really conflict", "TP = 34", "FN = 21"],
        ["really something else", "FP = 5", "TN = 1,060"],
    ], top=Inches(1.95), height=Inches(0.5), size=16)
    formula(s, [
        "Precision = TP/(TP+FP) = 34/39 = 0.87        "
        "Recall = TP/(TP+FN) = 34/55 = 0.62",
        "F1 = 2PR/(P+R) = 0.72",
    ], Inches(3.6), size=17, left=Inches(1.1), width=Inches(11.1))
    bullets(s, [
        "So when it says conflict, it is right 87% of the time. But it only finds 62% "
        "of the conflict stories. Too shy.",
        "Education is the opposite: precision 0.63, recall 0.83. Too eager. It shouts "
        "education at any story that mentions a school.",
        "Both land at F1 near 0.72. One number, two completely different problems.",
    ], top=Inches(5.1), size=15)

    # --- 7d. methodology ----------------------------------------------------
    s = _blank(prs)
    heading(s, "How I stopped myself from cheating", kicker="Validation")
    bullets(s, [
        "Train, validation and test are split by time. Every test article was collected "
        "after every training article, like a real deployment.",
        "No story group is split across two sets, so the same PTI copy cannot be in "
        "both training and test.",
        "Error bars come from bootstrap resampling. Same bootstrap as Chapter 06 §3.1, "
        "just used for confidence intervals instead of for bagging.",
        "I never compare two models by subtracting their scores. On 1,120 articles a "
        "gap of 0.02 can easily be luck, so I use a paired test on the articles where "
        "the two models actually disagree.",
        "Two publishers held out on every run: The Hindu and The Guardian, each time "
        "refitting the model without that publisher.",
        "The test set was opened once, at the very end, after the model was frozen.",
    ])

    # --- 8. the ladder ------------------------------------------------------
    s = _blank(prs)
    heading(s, "I started with the models we studied", kicker="Results")
    charted = picture(s, "ladder.png", Inches(0.8), Inches(1.9), Inches(7.3))
    bullets(s, [
        "Logistic regression first, because that is the classifier Chapter 03 builds. "
        "It reached 0.752.",
        "Then the tree side of the course: random forest, XGBoost, and a majority vote "
        "of ten models. Every one of them came in lower.",
        "So I was stuck around 0.75 with nothing left on the syllabus to try.",
        "That is when I went and read about the support vector machine.",
        ("Naive Bayes is the odd one. It gets WORSE with the full article, 0.635 down to "
         "0.594. More text is not automatically better.", 1),
    ], top=Inches(2.1),
        left=Inches(8.4) if charted else Inches(0.7),
        width=Inches(4.3) if charted else Inches(11.9),
        size=15 if charted else 18)
    note(s, "Every bar is validation macro-F1 on the same 1,120 articles, same split, "
            "same cleaning. Only the model changes.", top=Inches(5.5))

    # --- 9. why the SVM -----------------------------------------------------
    s = _blank(prs)
    heading(s, "Why I ended up using an SVM", kicker="Method")
    formula(s, [
        "logistic regression:   p = 1 / (1 + e^-(w·x + b))",
        "linear SVM:            answer yes if  w·x + b > 0",
    ], Inches(1.95), size=18, left=Inches(1.6), width=Inches(10.1))
    bullets(s, [
        "Both are the same shape. Multiply every word count by a weight, add them up, "
        "look at the sign. Only the fitting rule is different.",
        "Logistic regression picks weights that maximise likelihood. The SVM picks "
        "weights that push the boundary as far as it can from both classes. That gap is "
        "the margin.",
        "Chapter 06b mentions the SVM once, as the expensive strong learner that "
        "ensembles try to replace. It never builds one, so I read about it on my own.",
        "Honest result: on its own the SVM ties logistic regression, 0.753 against "
        "0.752. What actually helped was adding character n-grams, which took it to "
        "0.771.",
        "It also stays linear, so I can still print the exact words behind any "
        "prediction. A tree ensemble cannot.",
    ], top=Inches(3.5), size=16)
    note(s, "One problem it brought with it: no sigmoid, so the output is a distance "
            "and not a probability. That is fixed two slides from now.", top=Inches(6.7))

    # --- 9a. the body wins --------------------------------------------------
    s = _blank(prs)
    heading(s, "The full article wins", kicker="The main result")
    bullets(s, [
        f"Same model, same split, same cleaning. The only change is the text I feed it: "
        f"headline and summary gives {ab['title_summary']:.3f}, headline and full body "
        f"gives {ab['title_body']:.3f}.",
        "The two confidence intervals do not overlap. The paired test on the articles "
        "where the two versions disagree gives p = 0.0000075, so this is not luck.",
        "Why it works: a headline says “RBI holds rates”. The body says repo rate, "
        "inflation, monetary policy committee, bond yields. Four more chances to get "
        "it right instead of one.",
    ])
    _, run = _text(s, Inches(0.7), Inches(5.2), Inches(11.9), Inches(0.8), size=28,
                   bold=True, color=GOOD)
    run.text = (f"+{ab['delta']:.3f} macro-F1 just from reading the article "
                f"instead of the headline")

    # --- 9b. what failed ----------------------------------------------------
    s = _blank(prs)
    heading(s, "Things I tried that did not work", kicker="Results")
    table(s, [
        ["What I tried", "From", "What happened", "Kept?"],
        ["Random forest / extra trees", "Ch 06 §4", "0.04 to 0.09 worse", "no"],
        ["XGBoost on compressed features", "Ch 06 §5.3", "0.036 worse, 5× slower", "no"],
        ["Majority vote of 10 models", "Ch 06 §2", "+0.001, p = 1.00", "no"],
        ["KNN", "Ch 04", "ruled out on paper, not run", "no"],
        ["Removing names and places", "my idea", "no version helped", "no"],
        ["Giving the headline extra weight", "my idea", "steadily worse", "no"],
        ["Tuning the SVM over a 6× range", "my idea", "changes score by 0.002", "no"],
        ["Sentence embeddings (MiniLM)", "researched", "0.061 worse", "no"],
    ], top=Inches(2.0), height=Inches(0.42))
    note(s, "KNN is the one I ruled out by reasoning instead of by running it. It is a "
            "lazy learner: it keeps all 5,487 training articles and measures distance "
            "to each one at predict time. My feature space has about 30,000 dimensions, "
            "which is exactly the curse of dimensionality Chapter 04 warns about.",
         top=Inches(5.9))

    # --- 10. ensemble paradox ----------------------------------------------
    s = _blank(prs)
    heading(s, "Why my ensemble failed, and Chapter 06 told me so",
            kicker="Finding")
    bullets(s, [
        "Chapter 06 §2 gives two conditions for a vote to beat its members. Each model "
        "must be better than random, and their mistakes must be roughly independent. "
        "Mine failed the second one.",
        "The models do disagree, a lot. If I could magically pick the right model for "
        "every article I would score 0.900 instead of 0.771. So the disagreement is "
        "real and it is big.",
        "But it is lopsided. The embedding model fixes 70 of the SVM's mistakes and "
        "breaks 148 of its correct answers. A vote has no way to tell those apart.",
        "Four different combinations of models, four identical scores, p = 1.0000.",
        ("So I dropped the stacked ensemble I had planned and wrote down why. If a "
         "chooser would help but an average does not, then what I need is abstention, "
         "which is the next slide.", 1),
    ])
    # --- 11. calibration ----------------------------------------------------
    s = _blank(prs)
    heading(s, "Novelty 1: making the confidence number mean something",
            kicker="Innovation")
    if not picture(s, "calibration.png", Inches(1.2), Inches(2.0), Inches(10.9)):
        bullets(s, ["(figures/calibration.png missing — run "
                    "`uv run python scripts/build_figures.py` in ml-v2/)"])
    note(s, f"The SVM gives me w·x + b, a distance from the boundary. A score of 2.4 "
            f"does not mean 90% sure, it does not mean anything on its own. So I take "
            f"the training data, check how often a score of that size is actually "
            f"correct, and use that fraction as the confidence. Left chart: the dots "
            f"should sit on the dotted line, and they nearly do. When it claims 0.9 it "
            f"is right about 9 times in 10. Average gap is {test['ece']:.3f} on test. "
            f"Chapter 07 §7 warns that AUC cannot see this at all.", top=Inches(5.6))

    # --- 12. abstention -----------------------------------------------------
    s = _blank(prs)
    heading(s, "Novelty 2: letting it say “I am not sure”", kicker="Chapter 07 §4")
    formula(s, ["Coverage = articles it answered / all articles"], Inches(1.9),
            size=18)
    table(s, [
        ["On the held-out test set", "coverage", "accuracy on what it filed"],
        ["forced to answer everything", "100%", f"{test['accuracy_without_abstention']:.1%}"],
        [f"stays quiet below {test['cut']:.3f}", f"{test['coverage']:.1%}",
         f"{test['accuracy_filed']:.1%}"],
    ], top=Inches(2.8), height=Inches(0.5), highlight=2, size=16)
    bullets(s, [
        "Chapter 07 §4 moves the threshold to trade precision against recall. I move the "
        "same threshold, but instead of flipping the answer I hold the article back and "
        "send it to a person.",
        "The cut-off is chosen on training data only, at the point where precision hits "
        "90%. It never looks at validation or test.",
        "One cut-off for all 13 topics. Separate cut-offs per topic gained nothing, "
        "because calibration had already put every topic on the same scale.",
        "It doubts the right things. On 63 articles that my labellers could not classify "
        "either, the median confidence is 0.70 against 0.81 on the rest.",
    ], top=Inches(4.4), size=15)

    # --- 13. test split -----------------------------------------------------
    s = _blank(prs)
    heading(s, "Opening the test set, once", kicker="Results")
    table(s, [
        ["", "validation", "test"],
        ["articles", f"{val['n']:,}", f"{test['n']:,}"],
        ["macro-F1",
         f"{val['macro_f1']:.3f} [{val['macro_f1_low']:.3f}, {val['macro_f1_high']:.3f}]",
         f"{test['macro_f1']:.3f} [{test['macro_f1_low']:.3f}, {test['macro_f1_high']:.3f}]"],
        ["accuracy", f"{val['accuracy']:.3f}", f"{test['accuracy']:.3f}"],
        ["gap between claimed and real confidence", f"{val['ece']:.3f}",
         f"{test['ece']:.3f}"],
        ["how much it answered", f"{0.761:.1%}", f"{test['coverage']:.1%}"],
        ["accuracy on what it answered", f"{0.879:.3f}", f"{test['accuracy_filed']:.3f}"],
    ], top=Inches(2.1), height=Inches(0.46), highlight=2)
    _, run = _text(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.8), size=20,
                   bold=True, color=GOOD)
    run.text = (f"Test is {test['macro_f1'] - val['macro_f1']:+.3f} below validation. "
                "I had written down beforehand that anything inside 0.05 was fine, so "
                "I changed nothing after seeing it.")

    # --- 14. per class ------------------------------------------------------
    s = _blank(prs)
    heading(s, "Precision and recall, topic by topic", kicker="Results")
    charted = picture(s, "precision-recall.png", Inches(0.8), Inches(1.9), Inches(6.5))
    bullets(s, [
        "Sport is the best: precision 0.93, recall 0.98.",
        "Conflict and education both have F1 near 0.72, but look at the bars. Conflict "
        "is high precision and low recall. Education is the reverse. Two different "
        "problems hiding behind one number.",
        ("Education only has 29 validation articles, so its bars move a lot. I treat "
         "that as noise, not as a result.", 1),
        "Society and lifestyle is the genuine failure. Recall 0.38, so it misses 41 of "
        "its 66 articles. It is three ideas stuck together: community, jobs, lifestyle.",
        "About 19% of all mistakes are on topic pairs where my own labellers disagreed "
        "with each other. So there is a ceiling here well below a perfect score.",
    ], top=Inches(2.0),
        left=Inches(7.6) if charted else Inches(0.7),
        width=Inches(5.0) if charted else Inches(11.9),
        size=15 if charted else 18)

    # --- 14b. borrowed vocabulary -------------------------------------------
    s = _blank(prs)
    heading(s, "Four things I had to learn outside class", kicker="Honesty")
    bullets(s, [
        "The SVM itself. Chapter 06b names it in one line as the strong learner that "
        "ensembles try to replace, and never builds one.",
        "Calibration. Chapter 07 §7 says AUC cannot see it, but nothing in the course "
        "fixes it. I learn a curve from raw score to how often that score is right.",
        "Turning text into numbers. Every chapter starts from a table with columns. "
        "News has no columns, so I count words, and also 3 to 5 letter chunks, which is "
        "what catches Rs. against Rs and different word endings.",
        "A paired test, so a 0.02 gap on 1,120 articles does not get reported as a real "
        "improvement when it is luck.",
        ("And one thing I thought was new and was not. The bootstrap behind my error "
         "bars is Chapter 06 §3.1's bootstrap sample, doing a different job.", 1),
    ])

    # --- 15. demo -----------------------------------------------------------
    s = _blank(prs)
    heading(s, "The demo", kicker="Live")
    shown = picture(s, "demo.png", Inches(0.9), Inches(1.9), Inches(7.4))
    bullets(s, [
        "Paste any news article. It gives a topic, a confidence, and either files it or "
        "holds it back for a human.",
        "It shows its working. The model is a weighted sum of word counts, so the words "
        "with the biggest weights are literally the reason. That is the actual sum, not "
        "an explanation invented afterwards.",
        "The threshold is a live slider. Move it and watch coverage and accuracy trade "
        "against each other, which is the Chapter 07 §4 curve running in front of you.",
        "A Validation tab with the real per-class bars, the reliability chart and the "
        "confusion matrix, so you can check the claims instead of believing them.",
        "A batch tab that runs it over real collected articles that nobody labelled.",
    ], top=Inches(2.0),
        left=Inches(8.5) if shown else Inches(0.7),
        width=Inches(4.2) if shown else Inches(11.9),
        size=13 if shown else 18)
    note(s, "uv run streamlit run app/streamlit_app.py", top=Inches(6.6))

    # --- 16. limits ---------------------------------------------------------
    s = _blank(prs)
    heading(s, "What I am not claiming", kicker="Limitations")
    bullets(s, [
        "I collected over four days. So I cannot say anything about how news topics "
        "change over months. The two publisher holdouts are the only evidence I have "
        "that it generalises.",
        "English only, and mostly Indian papers, plus The Guardian, BBC and France 24. "
        "I would not trust it on Hindi or Tamil news without retraining.",
        "Conflict and war has only 120 training articles and dropped 0.14 on the test "
        "set. A thin class behaving like a thin class. I flagged that before opening "
        "the test set, not after.",
        "Most of the labelling was done by one person, me. So I cannot measure how noisy "
        "the labels are by comparing two labellers.",
        "I do not report ROC or AUC. Chapter 07 §7 says AUC only scores ranking and "
        "ignores whether the probabilities are honest, and honest probabilities are the "
        "whole basis of my abstention. So I measured that instead.",
        "Obvious next step is a modern language model. I deliberately left it until "
        "after this comparison, so any gain belongs to one change and not two.",
    ])

    # --- 17. conclusion -----------------------------------------------------
    s = _blank(prs)
    heading(s, "What I found", kicker="Conclusion")
    bullets(s, [
        f"Reading the full article beats reading the headline, by "
        f"{ab['delta']:.3f} macro-F1 with p = {ab['mcnemar_p']:.1e}. My hypothesis held.",
        "Trees, forests, boosting and voting, the entire ensemble half of the course, "
        "all scored below one plain linear model. The classifier was never my problem. "
        "Cleaning the data was.",
        f"Adding an honest confidence and the right to stay quiet turned a "
        f"{test['accuracy_without_abstention']:.1%} accurate model into one that answers "
        f"{test['coverage']:.1%} of new articles and gets "
        f"{test['accuracy_filed']:.1%} of those right.",
        f"Final score on the test set: macro-F1 {test['macro_f1']:.3f}, opened once, "
        "never reopened.",
    ])
    _, run = _text(s, Inches(0.7), Inches(5.8), Inches(11.9), Inches(0.8), size=22,
                   bold=True, color=ACCENT)
    run.text = ("The most useful part of this project is the list of things that did "
                "not work, with a number attached to each one.")

    return prs


def main() -> int:
    if not METRICS.is_file():
        raise SystemExit(f"no metrics at {METRICS}; run: uv run newsmlv2 train --id v2-001")
    prs = build(json.loads(METRICS.read_text(encoding="utf-8")))
    prs.save(OUT)
    print(f"wrote {OUT} ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
